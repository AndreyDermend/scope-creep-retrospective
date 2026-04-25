"""Dimitar the SCRUM Master — generates retrospective deck with QA loop."""

from __future__ import annotations

import multiprocessing as mp
import os
import subprocess
import time

from openai import OpenAI

from scope_creep.agents.base import Agent
from scope_creep.utils.code_extract import extract_python_code
from scope_creep.utils.transcript import Transcript, UIEvent, save_transcript


def qa_check_pptx(
    path: str,
    required_topics: list[str] | None = None,
    min_slides: int = 5,
    min_words_per_slide: int = 3,
) -> list[str]:
    """QA a .pptx file. Returns a list of issues (empty list = pass).

    Extracted from ScrumMaster so it can be unit-tested independently.
    """
    required_topics = required_topics or []
    issues: list[str] = []

    if not os.path.exists(path):
        return [f"{path} was not created"]

    try:
        from pptx import Presentation

        prs = Presentation(path)
    except Exception as e:  # noqa: BLE001
        return [f"Cannot open {path}: {e}"]

    slides = list(prs.slides)
    if len(slides) < min_slides:
        issues.append(f"Only {len(slides)} slides, need at least {min_slides}")

    all_text: list[str] = []
    for i, slide in enumerate(slides, 1):
        words: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                words.extend(shape.text_frame.text.split())
        all_text.extend(w.lower() for w in words)
        if len(words) < min_words_per_slide:
            issues.append(
                f"Slide {i} has fewer than {min_words_per_slide} words of content"
            )

    joined = " ".join(all_text)
    for topic in required_topics:
        if topic.lower() not in joined:
            issues.append(f"No slide mentions '{topic}'")

    return issues


class ScrumMaster(Agent):
    """Generate retrospective deck, then QA it by re-opening the .pptx.

    If any QA check fails, feed the specific issues back to the LLM and
    regenerate. Up to `max_attempts` times.
    """

    def __init__(
        self,
        name: str,
        system_prompt: str,
        inbox: mp.Queue,
        ui_queue: mp.Queue | None = None,
        target_file: str = "presentation.pptx",
        required_topics: list[str] | None = None,
        min_slides: int = 5,
        max_attempts: int = 3,
        model: str = "gpt-4.1-mini",
        temperature: float = 0.3,
    ) -> None:
        super().__init__(name, system_prompt, inbox, ui_queue, model, temperature)
        self.target_file = target_file
        self.required_topics = required_topics or ["scope", "lesson"]
        self.min_slides = min_slides
        self.max_attempts = max_attempts

    def run_qa(self) -> list[str]:
        return qa_check_pptx(
            self.target_file,
            required_topics=self.required_topics,
            min_slides=self.min_slides,
        )

    def run(self) -> None:  # pragma: no cover
        client = OpenAI()
        transcript = Transcript(agent=self.agent_name, model=self.model)
        messages: list[dict[str, str]] = [
            {"role": "system", "content": self.system_prompt}
        ]

        self.set_phase("waiting")
        self.emit("status", "Awaiting retrospective brief from lead")

        msg = self.inbox.get()
        briefing = msg.get("content", "")
        if briefing == Agent.task_done():
            return

        self.set_phase("design")
        self.emit("input", f"Received brief ({len(briefing)} chars)")

        initial_prompt = (
            briefing
            + "\n\nWrite Python code using python-pptx to create "
            f"'{self.target_file}'. Requirements:\n"
            f"- At least {self.min_slides} slides\n"
            "- Slide 1: title slide with project name\n"
            "- Content slides must cover: project overview, scope creep "
            "incident, what was removed, lessons learned, next steps\n"
            "- Every slide must have a visible title and body content\n"
            "- python-pptx is already installed; do not pip install\n"
            "- No main function, no __name__ guard, no argparse\n"
            "- Return ONLY a single ```python code block\n"
        )
        messages.append({"role": "user", "content": initial_prompt})

        success = False
        for attempt in range(1, self.max_attempts + 1):
            self.set_phase(f"attempt_{attempt}")
            self.emit("status", f"Generating deck (attempt {attempt})")

            try:
                resp = client.chat.completions.create(
                    model=self.model,
                    temperature=self.temperature,
                    messages=messages,
                )
                reply = resp.choices[0].message.content.strip()
            except Exception as e:  # noqa: BLE001
                self.emit("error", f"LLM error on attempt {attempt}: {e}")
                break

            messages.append({"role": "assistant", "content": reply})
            gen_code = extract_python_code(reply)

            if os.path.exists(self.target_file):
                os.remove(self.target_file)
            exec_result = subprocess.run(
                ["python", "-c", gen_code],
                capture_output=True,
                text=True,
                timeout=60,
            )
            time.sleep(1)

            issues = self.run_qa()
            if not issues:
                self.emit("qa", f"Attempt {attempt}: all checks pass ✓")
                self.emit("output", f"→ {self.target_file}")
                success = True
                break

            issue_lines = "\n".join(f"  ✗ {i}" for i in issues)
            self.emit(
                "qa", f"Attempt {attempt}: {len(issues)} issue(s)\n{issue_lines}"
            )

            if attempt < self.max_attempts:
                feedback = (
                    "QA found these problems with your deck:\n"
                    + "\n".join(f"- {i}" for i in issues)
                    + "\nFix them and resubmit. Return only a ```python code block."
                )
                if exec_result.stderr:
                    feedback += (
                        f"\n\nExecution stderr (may be relevant):\n"
                        f"{exec_result.stderr[:500]}"
                    )
                messages.append({"role": "user", "content": feedback})

        status = Agent.task_done() if success else "QA_FAILED"
        if self.outbox is not None:
            self.outbox.put({"from": self.agent_name, "content": status})
        self.emit("result" if success else "error", status)

        transcript.events = [UIEvent(**e) for e in self._events]
        transcript.messages = messages
        transcript.end_time = time.time()
        save_transcript(transcript, "transcripts")
