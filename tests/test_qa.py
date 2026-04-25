"""Tests for the PPTX QA checker.

Uses real python-pptx to generate decks with specific properties, then
runs them through `qa_check_pptx` to verify each failure mode is caught.
"""

from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches

from scope_creep.agents.scrum import qa_check_pptx


def _make_deck(path: Path, slides_spec: list[dict]) -> None:
    """Build a tiny deck. slides_spec = [{'title': str, 'body': str}, ...]"""
    prs = Presentation()
    blank = prs.slide_layouts[5]  # title only
    for spec in slides_spec:
        slide = prs.slides.add_slide(blank)
        # Title placeholder
        if slide.shapes.title is not None:
            slide.shapes.title.text = spec.get("title", "")
        # Add a textbox with body
        tb = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        tb.text_frame.text = spec.get("body", "")
    prs.save(str(path))


def test_valid_deck_passes(tmp_path):
    deck = tmp_path / "valid.pptx"
    _make_deck(deck, [
        {"title": "Title Slide", "body": "Customer churn prediction project"},
        {"title": "Overview", "body": "Scope of the project described"},
        {"title": "What happened", "body": "Scope creep incident occurred"},
        {"title": "Removed", "body": "Items that fell outside the scope"},
        {"title": "Lesson", "body": "Important lesson about churn analysis"},
        {"title": "Next steps", "body": "Keep the scope tight and lessons shared"},
    ])
    issues = qa_check_pptx(
        str(deck),
        required_topics=["scope", "lesson", "churn"],
        min_slides=5,
    )
    assert issues == []


def test_missing_file(tmp_path):
    issues = qa_check_pptx(
        str(tmp_path / "does_not_exist.pptx"),
        required_topics=["scope"],
        min_slides=5,
    )
    assert len(issues) == 1
    assert "was not created" in issues[0]


def test_too_few_slides(tmp_path):
    deck = tmp_path / "short.pptx"
    _make_deck(deck, [
        {"title": "Only", "body": "One slide here about scope and lesson and churn"},
    ])
    issues = qa_check_pptx(
        str(deck), required_topics=["scope", "lesson", "churn"], min_slides=5
    )
    assert any("Only 1 slides" in i for i in issues)


def test_missing_required_topic(tmp_path):
    deck = tmp_path / "no_topic.pptx"
    _make_deck(deck, [
        {"title": f"Slide {i}", "body": "filler " * 10}
        for i in range(5)
    ])
    issues = qa_check_pptx(
        str(deck),
        required_topics=["nonexistent_topic_xyz"],
        min_slides=5,
    )
    assert any("nonexistent_topic_xyz" in i for i in issues)


def test_thin_slide_detected(tmp_path):
    deck = tmp_path / "thin.pptx"
    _make_deck(deck, [
        {"title": "Full", "body": "Plenty of content here about scope lesson churn"},
        {"title": "Full", "body": "Plenty of content here about scope lesson churn"},
        {"title": "", "body": "hi"},  # too few words
        {"title": "Full", "body": "Plenty of content here about scope lesson churn"},
        {"title": "Full", "body": "Plenty of content here about scope lesson churn"},
    ])
    issues = qa_check_pptx(
        str(deck),
        required_topics=["scope", "lesson", "churn"],
        min_slides=5,
    )
    assert any("Slide 3" in i and "fewer than" in i for i in issues)
