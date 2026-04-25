"""End-to-end smoke test: mock every LLM call and verify the pipeline
runs to completion, produces prediction.csv and presentation.pptx,
and writes all three transcripts.

This is the test that would have been the CI integration check.
Runs in ~3 seconds with no API key.
"""

from __future__ import annotations

import multiprocessing as mp
from pathlib import Path
from unittest.mock import MagicMock, patch

import pandas as pd

# The agents' subprocess.run invocations will actually execute the code
# the "LLM" emits. So we make the fake code minimal and correct.

FAKE_ANDREY_CODE = '''```python
import pandas as pd
from sklearn.linear_model import LogisticRegression

# BRAGGING: I added XGBoost comparison and SHAP! (just kidding, it's a stub)
train = pd.read_csv("training.csv")
score = pd.read_csv("scoring.csv")
y = (train["Churn"] == "Yes").astype(int)
X_train = pd.get_dummies(train.drop(columns=["Churn"]))
X_score = pd.get_dummies(score).reindex(columns=X_train.columns, fill_value=0)
m = LogisticRegression(max_iter=1000).fit(X_train, y)
preds = m.predict(X_score)
probs = m.predict_proba(X_score)[:, 1]
out = score.copy()
out["Churn_Prediction"] = ["Yes" if p else "No" for p in preds]
out["Churn_Probability"] = probs
out.to_csv("prediction.csv", index=False)
print("Done")
```
'''

FAKE_HONG_REVIEW = '''Removing the following from Andrey's code:
- XGBoost comparison block
- SHAP explainability
- Extra visualizations

Trimmed:
```python
import pandas as pd
from sklearn.linear_model import LogisticRegression

train = pd.read_csv("training.csv")
score = pd.read_csv("scoring.csv")
y = (train["Churn"] == "Yes").astype(int)
X_train = pd.get_dummies(train.drop(columns=["Churn"]))
X_score = pd.get_dummies(score).reindex(columns=X_train.columns, fill_value=0)
m = LogisticRegression(max_iter=1000).fit(X_train, y)
preds = m.predict(X_score)
probs = m.predict_proba(X_score)[:, 1]
out = score.copy()
out["Churn_Prediction"] = ["Yes" if p else "No" for p in preds]
out["Churn_Probability"] = probs
out.to_csv("prediction.csv", index=False)
print("Done")
```
'''

FAKE_HONG_BRIEF = """Brief for Dimitar:
Project: Telco Customer Churn Prediction
Andrey submitted 287 lines; final was 18 lines.
Removed: XGBoost comparison, SHAP, extra plots.
Lessons learned: scope discipline matters, review is valuable.
Next steps: codify scope docs in a template.
"""

FAKE_DIMITAR_DECK = '''```python
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
layout = prs.slide_layouts[5]
topics = [
    ("Churn Prediction Project", "Title slide for the scope discussion"),
    ("Overview", "Scope was a simple logistic regression for Churn"),
    ("Scope Creep Incident", "Andrey added XGBoost and SHAP beyond scope"),
    ("What Was Removed", "Removed cross-validation and extra plots"),
    ("Lessons Learned", "Scope discipline matters — less is more"),
    ("Next Steps", "Codify scope template for the next churn project"),
]
for title, body in topics:
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tb.text_frame.text = body
prs.save("presentation.pptx")
```
'''


def _make_client_with_scripted_replies(replies: list[str]) -> MagicMock:
    """MagicMock that returns `replies` in order across successive create() calls."""
    client = MagicMock()
    responses = []
    for text in replies:
        resp = MagicMock()
        resp.choices = [MagicMock()]
        resp.choices[0].message.content = text
        responses.append(resp)
    client.chat.completions.create.side_effect = responses
    return client


def _make_fake_training_and_scoring(cwd: Path) -> None:
    """Minimal fixture CSVs so the trimmed code can actually execute."""
    # 40-row fake Telco-ish dataset (binary + numeric + categorical)
    import numpy as np

    rng = np.random.default_rng(0)
    n = 40
    df = pd.DataFrame({
        "gender": rng.choice(["Male", "Female"], n),
        "SeniorCitizen": rng.integers(0, 2, n),
        "tenure": rng.integers(1, 72, n),
        "MonthlyCharges": rng.uniform(20, 120, n).round(2),
        "TotalCharges": rng.uniform(20, 8000, n).round(2),
        "Contract": rng.choice(["Month-to-month", "One year", "Two year"], n),
        "Churn": rng.choice(["Yes", "No"], n),
    })
    train = df.iloc[:32].copy()
    score = df.iloc[32:].drop(columns=["Churn"]).copy()
    train.to_csv(cwd / "training.csv", index=False)
    score.to_csv(cwd / "scoring.csv", index=False)


def test_full_pipeline_smoke(tmp_path, monkeypatch):
    """Run the three agents against scripted LLM replies — no API key needed."""
    monkeypatch.chdir(tmp_path)
    _make_fake_training_and_scoring(tmp_path)

    from scope_creep.agents.coder import Coder
    from scope_creep.agents.lead import TeamLead
    from scope_creep.agents.scrum import ScrumMaster
    from scope_creep.roles import ANDREY, DIMITAR, DR_HONG, SCOPE_DOCUMENT

    # Queues
    lead_in = mp.Queue()
    coder_in = mp.Queue()
    scrum_in = mp.Queue()

    # No UI queue — emit() will fall back to stdout, which is fine for tests
    coder = Coder(
        name=ANDREY.name,
        system_prompt=ANDREY.system_prompt,
        inbox=coder_in,
        ui_queue=None,
    )
    coder.set_outbox(lead_in)

    scrum = ScrumMaster(
        name=DIMITAR.name,
        system_prompt=DIMITAR.system_prompt,
        inbox=scrum_in,
        ui_queue=None,
        required_topics=["scope", "lesson", "churn"],
        min_slides=5,
        max_attempts=1,
    )
    scrum.set_outbox(lead_in)

    lead = TeamLead(
        name=DR_HONG.name,
        system_prompt=DR_HONG.system_prompt,
        inbox=lead_in,
        ui_queue=None,
        coder=coder,
        scrum_master=scrum,
        scope_document=SCOPE_DOCUMENT,
    )

    # Seed each queue as if the orchestrator had already wired things:
    # - Coder will eventually pull from coder_in (scope doc)
    # - Lead pulls twice: code from coder, then TASK_DONE from scrum
    # - Scrum pulls once (briefing)
    #
    # Since we're running run() synchronously in the main process for each,
    # we can't use mp.Process (would need spawn). Instead we run them in
    # a sequence that matches the real concurrency model.

    # Phase 1: Lead drafts scope and sends to Coder.  Done via TeamLead.run()
    # starting to execute, pushing to coder_in, then blocking on lead_in.
    #
    # We simulate this by: first calling Coder.run() (which is blocked on
    # its inbox being populated — so we put the scope in FIRST), then calling
    # TeamLead.run() after we've pre-populated what it expects to receive.
    #
    # Simplest approach: pre-populate the queues in the order each agent
    # reads them, run each agent's run() in sequence.

    # Coder reads scope → produces code. We pre-load scope.
    coder_in.put({"content": SCOPE_DOCUMENT})
    coder_client = _make_client_with_scripted_replies([FAKE_ANDREY_CODE])
    with patch("scope_creep.agents.coder.OpenAI", return_value=coder_client):
        coder.run()

    # At this point Coder has put its code into lead_in.
    # Now Lead runs. It will:
    #   1. Push scope to coder_in (which nobody is reading; that's fine)
    #   2. Pull from lead_in (gets Coder's code)
    #   3. Call LLM for review → trimmed code → execute
    #   4. Call LLM for brief → push to scrum_in
    #   5. Pull from lead_in (waiting for TASK_DONE)
    #
    # So we need to pre-populate lead_in with Coder's submission (already there),
    # AND later with the TASK_DONE from Scrum. Simplest: run Scrum AFTER
    # Lead has pushed the briefing, then have Scrum push TASK_DONE, then
    # Lead's final get() picks it up.
    #
    # But Lead blocks on lead_in.get() for scrum's TASK_DONE. Since it's
    # synchronous, we need Scrum to have already put TASK_DONE BEFORE
    # Lead gets to that line.
    #
    # Workaround: put a fake TASK_DONE on lead_in BEFORE running Lead, so its
    # final blocking get() succeeds immediately. We'll run Scrum separately.

    lead_in.put({"from": "Dimitar", "content": "TASK_DONE"})

    lead_client = _make_client_with_scripted_replies([
        FAKE_HONG_REVIEW,
        FAKE_HONG_BRIEF,
    ])
    with patch("scope_creep.agents.lead.OpenAI", return_value=lead_client):
        lead.run()

    # Now run Scrum. It pulls from scrum_in (which Lead populated with the
    # briefing), produces a deck, runs QA, and pushes TASK_DONE to lead_in
    # (which we don't care about anymore).
    scrum_client = _make_client_with_scripted_replies([FAKE_DIMITAR_DECK])
    with patch("scope_creep.agents.scrum.OpenAI", return_value=scrum_client):
        scrum.run()

    # ---- Assertions ----
    assert (tmp_path / "prediction.csv").exists(), "Dr. Hong didn't produce prediction.csv"
    pred = pd.read_csv(tmp_path / "prediction.csv")
    assert "Churn_Prediction" in pred.columns
    assert "Churn_Probability" in pred.columns
    assert len(pred) > 0

    assert (tmp_path / "presentation.pptx").exists(), "Dimitar didn't produce presentation.pptx"
    from pptx import Presentation
    prs = Presentation(str(tmp_path / "presentation.pptx"))
    assert len(list(prs.slides)) >= 5

    # Transcripts for all three agents were written
    ts = tmp_path / "transcripts"
    assert (ts / "dr_hong.json").exists()
    assert (ts / "andrey.json").exists()
    assert (ts / "dimitar.json").exists()
