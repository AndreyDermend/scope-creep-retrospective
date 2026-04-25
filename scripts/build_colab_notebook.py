"""Build the Colab submission notebook from the package source."""
from __future__ import annotations

import json
from pathlib import Path

ROOT = Path("/home/claude/scope-creep-retrospective")


def read(rel: str) -> str:
    return (ROOT / rel).read_text()


def md(text: str) -> dict:
    return {"cell_type": "markdown", "metadata": {}, "source": text}


def code(src: str) -> dict:
    return {
        "cell_type": "code",
        "execution_count": None,
        "metadata": {},
        "outputs": [],
        "source": src,
    }


cells: list[dict] = []

# ============================================================
# Title
# ============================================================
cells.append(md(
    "# Scope-Creep Retrospective — Three-Agent ML Pipeline\n"
    "\n"
    "A three-agent OpenAI pipeline where the agents **disagree by design**:\n"
    "\n"
    "- **Dr. Hong** (Team Lead) — writes the scope, reviews submitted code, "
    "trims out-of-scope additions, executes the final version.\n"
    "- **Andrey Dermendzhiev** (Coder) — fulfills the scope then over-delivers. "
    "Extra models, cross-validation, SHAP, bonus plots. Does *not* execute.\n"
    "- **Dimitar Dermendzhiev** (SCRUM Master) — builds a \"lessons learned\" "
    "deck, then re-opens the `.pptx` to run a structural QA loop. Regenerates "
    "up to 3 times if anything fails.\n"
    "\n"
    "**Dataset:** [Telco Customer Churn](https://github.com/IBM/"
    "telco-customer-churn-on-icp4d) — 7,043 rows, binary `Churn` target. "
    "Downloaded automatically below.\n"
    "\n"
    "Based on the multi-agent multiprocessing pattern from "
    "`5_heart_attack_analysis_v2.ipynb`, extended with scope-creep narrative, "
    "self-QA loops, structured event logging, and a live Rich terminal UI.\n"
    "\n"
    "## Flow\n"
    "\n"
    "```\n"
    "Dr. Hong → Andrey (scope)\n"
    "Andrey → Dr. Hong (inflated code, ~280 lines)\n"
    "Dr. Hong reviews, trims (~40 lines), executes → prediction.csv\n"
    "Dr. Hong → Dimitar (retrospective brief)\n"
    "Dimitar generates deck, runs QA; regenerates if needed → presentation.pptx\n"
    "```"
))

# ============================================================
# Install
# ============================================================
cells.append(md("## 1. Install dependencies"))
cells.append(code(
    "!pip install -q openai pandas scikit-learn python-pptx rich"
))

# ============================================================
# Dataset prep
# ============================================================
cells.append(md(
    "## 2. Download and split the dataset\n"
    "\n"
    "Telco Churn → `training.csv` (with label) + `scoring.csv` (label stripped)."
))
cells.append(code(read("src/scope_creep/data_prep.py")))
cells.append(code("prepare_data()\nprint('Data ready.')"))

# ============================================================
# Utilities
# ============================================================
cells.append(md("## 3. Utilities — code extraction and transcripts"))
cells.append(code(read("src/scope_creep/utils/code_extract.py")))
cells.append(code(read("src/scope_creep/utils/transcript.py")))

# ============================================================
# Agents
# ============================================================
cells.append(md(
    "## 4. The three agents\n"
    "\n"
    "Each agent is an `mp.Process` with its own OpenAI client. They communicate "
    "through `mp.Queue` objects. The `emit()` method publishes structured "
    "events to a shared UI queue and also persists them to the transcript."
))

# Base — strip the import that came from utils.transcript (we inlined it)
base_src = read("src/scope_creep/agents/base.py")
base_src = base_src.replace(
    "from scope_creep.utils.transcript import UIEvent\n", ""
)
cells.append(code(base_src))

# Coder
coder_src = read("src/scope_creep/agents/coder.py")
coder_src = coder_src.replace(
    "from scope_creep.agents.base import Agent\n"
    "from scope_creep.utils.code_extract import extract_python_code, strip_code_fences\n"
    "from scope_creep.utils.transcript import Transcript, UIEvent, save_transcript\n",
    "",
)
cells.append(code(coder_src))

# Scrum
scrum_src = read("src/scope_creep/agents/scrum.py")
scrum_src = scrum_src.replace(
    "from scope_creep.agents.base import Agent\n"
    "from scope_creep.utils.code_extract import extract_python_code\n"
    "from scope_creep.utils.transcript import Transcript, UIEvent, save_transcript\n",
    "",
)
cells.append(code(scrum_src))

# Lead
lead_src = read("src/scope_creep/agents/lead.py")
lead_src = lead_src.replace(
    "from scope_creep.agents.base import Agent\n"
    "from scope_creep.utils.code_extract import extract_python_code\n"
    "from scope_creep.utils.transcript import Transcript, UIEvent, save_transcript\n",
    "",
)
# lead.py has a second import of strip_code_fences inline — keep it module-scope
lead_src = lead_src.replace(
    "        from scope_creep.utils.code_extract import strip_code_fences\n\n",
    "",
)
cells.append(code(lead_src))

# ============================================================
# Roles
# ============================================================
cells.append(md("## 5. Roles and scope document"))
cells.append(code(read("src/scope_creep/roles.py")))

# ============================================================
# UI
# ============================================================
cells.append(md(
    "## 6. Live terminal UI\n"
    "\n"
    "Three side-by-side Rich panels showing each agent's status, thinking, "
    "and outputs in real time. Colab renders Rich output inline so this "
    "works in the notebook too."
))
live_src = read("src/scope_creep/ui/live.py")
cells.append(code(live_src))

# ============================================================
# Run
# ============================================================
cells.append(md(
    "## 7. Run the pipeline\n"
    "\n"
    "⚠️ You'll need an OpenAI API key. Paste it when prompted. A full run "
    "costs ~$0.05 on `gpt-4.1-mini` and takes 60-90 seconds."
))
cells.append(code(
    "import os\n"
    "import multiprocessing as mp\n"
    "from getpass import getpass\n"
    "\n"
    "if 'OPENAI_API_KEY' not in os.environ:\n"
    "    os.environ['OPENAI_API_KEY'] = getpass('Enter your OpenAI API key: ')\n"
    "\n"
    "# Queues\n"
    "lead_in = mp.Queue()\n"
    "coder_in = mp.Queue()\n"
    "scrum_in = mp.Queue()\n"
    "ui_queue = mp.Queue()\n"
    "\n"
    "# Agents\n"
    "coder = Coder(name=ANDREY.name, system_prompt=ANDREY.system_prompt,\n"
    "              inbox=coder_in, ui_queue=ui_queue, temperature=0.5)\n"
    "coder.set_outbox(lead_in)\n"
    "\n"
    "scrum = ScrumMaster(name=DIMITAR.name, system_prompt=DIMITAR.system_prompt,\n"
    "                    inbox=scrum_in, ui_queue=ui_queue,\n"
    "                    required_topics=['scope', 'lesson', 'churn'],\n"
    "                    min_slides=5, max_attempts=3)\n"
    "scrum.set_outbox(lead_in)\n"
    "\n"
    "lead = TeamLead(name=DR_HONG.name, system_prompt=DR_HONG.system_prompt,\n"
    "                inbox=lead_in, ui_queue=ui_queue,\n"
    "                coder=coder, scrum_master=scrum,\n"
    "                scope_document=SCOPE_DOCUMENT)\n"
    "\n"
    "# Launch\n"
    "coder.start()\n"
    "scrum.start()\n"
    "lead.start()\n"
    "\n"
    "# Drive the UI in the main process\n"
    "ui = LiveUI([DR_HONG.name, ANDREY.name, DIMITAR.name], ui_queue)\n"
    "ui.consume_until_done(timeout=600)\n"
    "\n"
    "coder.join(timeout=30)\n"
    "lead.join(timeout=30)\n"
    "scrum.join(timeout=30)\n"
    "print('\\nDone.')"
))

# ============================================================
# Verify outputs
# ============================================================
cells.append(md(
    "## 8. Inspect the outputs\n"
    "\n"
    "After the run you should have `prediction.csv`, `presentation.pptx`, "
    "and per-agent transcripts in `transcripts/`."
))
cells.append(code(
    "import pandas as pd\n"
    "from pptx import Presentation\n"
    "\n"
    "print('=== prediction.csv ===')\n"
    "pred = pd.read_csv('prediction.csv')\n"
    "print(f'{len(pred):,} rows, predicted churn rate '\n"
    "      f'{(pred[\"Churn_Prediction\"] == \"Yes\").mean():.1%}')\n"
    "print(pred.head())\n"
    "\n"
    "print('\\n=== presentation.pptx ===')\n"
    "prs = Presentation('presentation.pptx')\n"
    "for i, slide in enumerate(prs.slides, 1):\n"
    "    title = ''\n"
    "    for shape in slide.shapes:\n"
    "        if shape.has_text_frame and shape.text_frame.text.strip():\n"
    "            title = shape.text_frame.text.strip().split('\\n')[0]\n"
    "            break\n"
    "    print(f'  Slide {i}: {title}')"
))

# ============================================================
# Download artifacts
# ============================================================
cells.append(md(
    "## 9. Download the artifacts (Colab only)\n"
    "\n"
    "Run this cell to download `prediction.csv` and `presentation.pptx` to "
    "your local machine."
))
cells.append(code(
    "try:\n"
    "    from google.colab import files\n"
    "    files.download('prediction.csv')\n"
    "    files.download('presentation.pptx')\n"
    "except ImportError:\n"
    "    print('Not in Colab — files are in the current working directory.')"
))

# ============================================================
# Assemble
# ============================================================
notebook = {
    "nbformat": 4,
    "nbformat_minor": 0,
    "metadata": {
        "colab": {"provenance": []},
        "kernelspec": {"name": "python3", "display_name": "Python 3"},
        "language_info": {"name": "python"},
    },
    "cells": cells,
}

out = ROOT / "notebooks" / "scope_creep_colab.ipynb"
out.parent.mkdir(parents=True, exist_ok=True)
out.write_text(json.dumps(notebook, indent=1))
print(f"Wrote {out}")
print(f"Cells: {len(cells)}")
